import base64
import io
import re
import zipfile
from typing import Any

try:
    import pikepdf
    HAS_PIKEPDF = True
except Exception:
    HAS_PIKEPDF = False

try:
    from oletools.olevba import VBA_Parser
    HAS_OLEVBA = True
except Exception:
    HAS_OLEVBA = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except Exception:
    HAS_DOCX = False

try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except Exception:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except Exception:
    HAS_PPTX = False

from app.api.schemas import DocumentScanResponse, DocumentThreat

URL_RE = re.compile(r"https?://[^\s\"<>]+", re.IGNORECASE)

def extract_urls_from_text(text: str) -> list[str]:
    if not text: return []
    return list(set(URL_RE.findall(text)))

def check_pdf_for_threats(file_bytes: bytes) -> dict[str, Any]:
    threats = []
    has_js = False
    links = []
    
    if not HAS_PIKEPDF:
        return {"threats": [DocumentThreat(type="warning", description="pikepdf not installed; limited PDF analysis", severity="low")], "has_javascript": False, "links": []}
    
    try:
        # pikepdf gracefully handles obfuscated streams
        with pikepdf.Pdf.open(io.BytesIO(file_bytes)) as pdf:
            if hasattr(pdf, "Root") and "/Names" in pdf.Root and "/JavaScript" in pdf.Root.Names:
                has_js = True
                threats.append(DocumentThreat(type="javascript", description="Document contains embedded JavaScript", severity="high"))
            
            if hasattr(pdf, "Root") and "/OpenAction" in pdf.Root:
                has_js = True # often used for JS or Launch
                threats.append(DocumentThreat(type="javascript", description="Document contains an OpenAction trigger", severity="critical"))
                
            for page in pdf.pages:
                if "/Annots" in page:
                    for annot in page.Annots:
                        if str(annot.get("/Subtype", "")) == "/Link" and "/A" in annot:
                            action = annot["/A"]
                            if str(action.get("/S", "")) == "/URI":
                                uri = str(action.get("/URI", ""))
                                if uri:
                                    links.append(uri)
                            elif str(action.get("/S", "")) == "/Launch":
                                threats.append(DocumentThreat(type="embedded_file", description="Page contains a /Launch action (often malware dropper)", severity="critical"))
                                
                if "/AA" in page:
                    threats.append(DocumentThreat(type="javascript", description="Page contains automatic action triggers (/AA)", severity="critical"))
                    has_js = True

    except Exception as exc:
        threats.append(DocumentThreat(type="parse_error", description=f"PDF parsing failed: {exc}", severity="low"))
    
    return {"threats": threats, "has_macros": False, "has_javascript": has_js, "links": list(set(links))}

def analyze_macros(file_bytes: bytes) -> list[DocumentThreat]:
    threats = []
    if not HAS_OLEVBA:
        threats.append(DocumentThreat(type="warning", description="oletools not installed; macro deep inspection disabled", severity="low"))
        return threats

    try:
        vbaparser = VBA_Parser("in_memory", data=file_bytes)
        if vbaparser.detect_vba_macros():
            threats.append(DocumentThreat(type="macro", description="Document contains VBA macros", severity="high"))
            for (filename, stream_path, vba_filename, vba_code) in vbaparser.extract_macros():
                code_lower = vba_code.lower()
                if "autoopen" in code_lower or "autoexec" in code_lower or "document_open" in code_lower:
                    threats.append(DocumentThreat(type="macro", description="Macro executes automatically on open", severity="critical"))
                if "shell" in code_lower or "createobject(\"wscript.shell\")" in code_lower.replace(" ", ""):
                    threats.append(DocumentThreat(type="macro", description="Macro attempts to execute shell commands", severity="critical"))
                if "chr(" in code_lower or "base64" in code_lower:
                    threats.append(DocumentThreat(type="macro", description="Macro contains obfuscated strings (Chr/Base64)", severity="high"))
        vbaparser.close()
    except Exception as exc:
        threats.append(DocumentThreat(type="parse_error", description=f"VBA Macro parsing failed: {exc}", severity="low"))
    return threats

def parse_docx(file_bytes: bytes) -> dict[str, Any]:
    threats = analyze_macros(file_bytes)
    links = []
    has_macros = any(t.type == "macro" for t in threats)
    
    if HAS_DOCX:
        try:
            doc = DocxDocument(io.BytesIO(file_bytes))
            for para in doc.paragraphs:
                links.extend(extract_urls_from_text(para.text))
            for rel in doc.part.rels.values():
                if "hyperlink" in rel.reltype:
                    links.append(rel.target_ref)
        except Exception as exc:
            threats.append(DocumentThreat(type="parse_error", description=f"DOCX text parsing failed: {exc}", severity="low"))
    
    return {"threats": threats, "has_macros": has_macros, "has_javascript": False, "links": list(set(links))}

def parse_xlsx(file_bytes: bytes) -> dict[str, Any]:
    threats = analyze_macros(file_bytes)
    links = []
    has_macros = any(t.type == "macro" for t in threats)
    
    if HAS_OPENPYXL:
        try:
            wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if isinstance(cell, str):
                            links.extend(extract_urls_from_text(cell))
        except Exception as exc:
            threats.append(DocumentThreat(type="parse_error", description=f"XLSX text parsing failed: {exc}", severity="low"))
    
    return {"threats": threats, "has_macros": has_macros, "has_javascript": False, "links": list(set(links))}

def parse_pptx(file_bytes: bytes) -> dict[str, Any]:
    threats = analyze_macros(file_bytes)
    links = []
    has_macros = any(t.type == "macro" for t in threats)
    
    if HAS_PPTX:
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        links.extend(extract_urls_from_text(shape.text))
                    if hasattr(shape, "has_hyperlink") and shape.has_hyperlink:
                        links.append(shape.hyperlink.address or "")
        except Exception as exc:
            threats.append(DocumentThreat(type="parse_error", description=f"PPTX text parsing failed: {exc}", severity="low"))
    
    return {"threats": threats, "has_macros": has_macros, "has_javascript": False, "links": list(set(links))}

def inspect_document(filename: str, content_base64: str, mime_type: str) -> DocumentScanResponse:
    if "," in content_base64:
        content_base64 = content_base64.split(",", 1)[1]
    file_bytes = base64.b64decode(content_base64)
    ext = filename.lower().split(".")[-1] if "." in filename else ""
    
    if mime_type == "application/pdf" or ext == "pdf":
        result = check_pdf_for_threats(file_bytes)
    elif mime_type in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword") or ext in ("docx", "doc", "docm"):
        result = parse_docx(file_bytes)
    elif mime_type in ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel") or ext in ("xlsx", "xls", "xlsm"):
        result = parse_xlsx(file_bytes)
    elif mime_type in ("application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint") or ext in ("pptx", "ppt", "pptm"):
        result = parse_pptx(file_bytes)
    else:
        # Fallback to pure macro check for unknown files just in case
        threats = analyze_macros(file_bytes)
        has_macros = any(t.type == "macro" for t in threats)
        if not threats:
            threats.append(DocumentThreat(type="unsupported", description=f"Unsupported file type: {mime_type}", severity="low"))
        return DocumentScanResponse(
            filename=filename,
            mime_type=mime_type,
            threats_found=threats,
            has_macros=has_macros,
            has_javascript=False,
            external_links=[],
            risk_score=100 if has_macros else 0,
            prediction="malicious" if has_macros else "unknown"
        )
    
    risk = 0
    # Add logic to avoid duplicate scoring if multiple threats are found, but a simple sum works fine
    seen_critical = False
    for t in result["threats"]:
        if t.severity == "critical":
            risk += 50
            seen_critical = True
        elif t.severity == "high":
            risk += 25
        elif t.severity == "medium":
            risk += 10
        else:
            risk += 2
            
    # Guarantee critical threats push score over malicious threshold (70+)
    if seen_critical and risk < 75:
        risk = 75
    
    risk = min(risk, 100)
    prediction = "clean" if risk < 15 else "suspicious" if risk < 50 else "malicious"
    
    return DocumentScanResponse(
        filename=filename,
        mime_type=mime_type,
        threats_found=result["threats"],
        has_macros=result["has_macros"],
        has_javascript=result.get("has_javascript", False),
        external_links=result["links"],
        risk_score=risk,
        prediction=prediction
    )
